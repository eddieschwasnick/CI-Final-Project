"""
build_presentation.py
---------------------
Generate reports/presentation.pptx for the 10-minute solo project presentation.

Run:  python reports/build_presentation.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG_DIR  = os.path.join(BASE_DIR, "reports", "figures")
OUT_PATH = os.path.join(BASE_DIR, "reports", "presentation.pptx")

NAVY  = RGBColor(0x1F, 0x3A, 0x5F)
GREY  = RGBColor(0x55, 0x55, 0x55)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED   = RGBColor(0xC0, 0x39, 0x2B)


def add_title_box(slide, text, left, top, width, height,
                  size=28, bold=True, color=NAVY, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def add_body_box(slide, bullets, left, top, width, height,
                 size=16, color=BLACK, bullet_char="•"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        prefix = f"{bullet_char}  " if bullet_char else ""
        run.text = prefix + line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.space_after = Pt(8)
    return tf


def add_footer(slide, slide_num, total):
    left = Inches(0.3)
    top  = Inches(7.1)
    add_title_box(slide, f"Schwasnick  |  STAT 6990  |  {slide_num}/{total}",
                  left, top, Inches(9.5), Inches(0.3),
                  size=10, bold=False, color=GREY)


def add_image(slide, image_path, left, top, width=None, height=None):
    return slide.shapes.add_picture(image_path, left, top,
                                    width=width, height=height)


def main():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]
    total_slides = 10

    # ── Slide 1: Title ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s,
        "Does the Information You Give an LLM\nChange the Portfolio It Builds?",
        Inches(0.5), Inches(1.6), Inches(9), Inches(1.8),
        size=32, color=NAVY, align=PP_ALIGN.CENTER)
    add_title_box(s,
        "A Randomized Experiment on Information Conditioning\n"
        "and LLM Allocation Behavior",
        Inches(0.5), Inches(3.4), Inches(9), Inches(1.0),
        size=18, bold=False, color=GREY, align=PP_ALIGN.CENTER)
    add_title_box(s, "Edward Schwasnick",
        Inches(0.5), Inches(5.5), Inches(9), Inches(0.4),
        size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_title_box(s, "STAT 6990  •  Causal Inference  •  Spring 2026",
        Inches(0.5), Inches(5.95), Inches(9), Inches(0.4),
        size=14, bold=False, color=GREY, align=PP_ALIGN.CENTER)
    add_title_box(s, "April 29, 2026",
        Inches(0.5), Inches(6.4), Inches(9), Inches(0.4),
        size=14, bold=False, color=GREY, align=PP_ALIGN.CENTER)

    # ── Slide 2: Research Question ────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Research Question",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_body_box(s, [
        "LLMs increasingly produce structured financial recommendations.",
        "Vendors claim their AI analysts give better advice because of better data.",
        "That's a causal claim — one worth testing.",
    ], Inches(0.6), Inches(1.5), Inches(8.8), Inches(1.8), size=17)

    # Highlighted question box
    add_title_box(s,
        "Holding model, temperature, universe, and instructions fixed,\n"
        "does the information in the prompt causally change\n"
        "the portfolio an LLM constructs?",
        Inches(0.8), Inches(3.6), Inches(8.4), Inches(1.8),
        size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_body_box(s, [
        "Approach: randomized 4-arm experiment where the only thing that varies "
        "across arms is the data block in the prompt.",
    ], Inches(0.6), Inches(5.8), Inches(8.8), Inches(1.0), size=16, color=GREY)
    add_footer(s, 2, total_slides)

    # ── Slide 3: Experimental Design ──────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Experimental Design",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))

    add_body_box(s, [
        "Stock universe: 20 U.S. large-caps across 6 sectors",
        "Model: claude-sonnet-4-6  •  Temperature: 0.7  •  Seed: 42",
        "30 runs per condition  =  120 API calls  •  randomized trial order",
        "Output: JSON with weights for all 20 tickers, summing to 100",
    ], Inches(0.6), Inches(1.3), Inches(8.8), Inches(1.8), size=16)

    # Four-arm table
    conditions = [
        ("Control",     "ticker list only",                  "#6c757d"),
        ("Fundamental", "+ fundamentals table\n(P/E, growth, leverage)", "#1f77b4"),
        ("Technical",   "+ technicals table\n(returns, beta, vol, MA)",   "#2ca02c"),
        ("Combined",    "+ both tables",                     "#d62728"),
    ]
    col_w = Inches(2.2); gap = Inches(0.1); left_start = Inches(0.55)
    top   = Inches(3.5); height = Inches(2.6)
    for i, (name, desc, hexcolor) in enumerate(conditions):
        left = left_start + i * (col_w + gap)
        box = s.shapes.add_textbox(left, top, col_w, height)
        tf = box.text_frame
        tf.word_wrap = True
        # title cell
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = name
        r1.font.size = Pt(18); r1.font.bold = True
        hx = hexcolor.lstrip("#")
        r1.font.color.rgb = RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        # body
        for line in desc.split("\n"):
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = line
            r.font.size = Pt(13); r.font.color.rgb = BLACK
            p.space_before = Pt(6)

    add_body_box(s, [
        "Estimand: ATE of each arm vs. control on 5 portfolio outcomes.",
    ], Inches(0.6), Inches(6.4), Inches(8.8), Inches(0.6), size=15, color=GREY)
    add_footer(s, 3, total_slides)

    # ── Slide 4: Outcomes & Estimation ────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Outcomes and Estimation",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_body_box(s, [
        "HHI  =  Σ wᵢ²                             (stock-level concentration)",
        "Sector HHI  =  Σ ( Σ_{i∈sector} wᵢ )²       (sector-level concentration)",
        "Portfolio β  =  Σ wᵢ · βᵢ",
        "Portfolio vol  =  Σ wᵢ · σᵢ",
        "Breadth  =  count of wᵢ > 0",
    ], Inches(0.8), Inches(1.5), Inches(8.6), Inches(2.6), size=16)

    add_title_box(s, "Inference",
        Inches(0.5), Inches(4.3), Inches(9), Inches(0.5),
        size=20, bold=True, color=NAVY)
    add_body_box(s, [
        "Primary:  Welch's t-test on diff-in-means vs. control.",
        "Regression:  Yᵢ = α + β₁·Fund + β₂·Tech + β₃·Comb + εᵢ   (control omitted).",
        "Placebo #1:  5,000 permutations of condition labels.",
        "Placebo #2:  1,000 random splits within the control group.",
        "Sensitivity:  re-estimate after trimming top/bottom HHI per condition.",
    ], Inches(0.8), Inches(4.9), Inches(8.6), Inches(2.3), size=15)
    add_footer(s, 4, total_slides)

    # ── Slide 5: Headline Results ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Headline: Prompt Information Moves Portfolios",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    # Table: means by condition
    rows, cols = 6, 5
    table_shape = s.shapes.add_table(rows, cols,
        Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.2)).table
    headers = ["", "Control", "Fundamental", "Technical", "Combined"]
    for c, h in enumerate(headers):
        cell = table_shape.cell(0, c)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.CENTER
            for run in para.runs:
                run.font.bold = True; run.font.size = Pt(14); run.font.color.rgb = WHITE
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY

    body_rows = [
        ("HHI",               "0.060", "0.103", "0.076", "0.059"),
        ("Sector HHI",        "0.209", "0.366", "0.173", "0.182"),
        ("Portfolio β",       "0.92",  "1.05",  "0.78",  "0.85"),
        ("Portfolio vol (%)", "29.9",  "30.3",  "27.7",  "28.5"),
        ("Breadth",           "20.0",  "20.0",  "20.0",  "20.0"),
    ]
    for r, row in enumerate(body_rows, start=1):
        for c, val in enumerate(row):
            cell = table_shape.cell(r, c)
            cell.text = val
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(13)
                    if c == 0:
                        run.font.bold = True

    add_body_box(s, [
        "Fundamental prompt  →  most concentrated,  highest β (Tech-heavy).",
        "Technical prompt  →  most diversified,  lowest β  (momentum + defensive).",
        "Combined  ≠  average of the two; HHI actually drops vs. control.",
    ], Inches(0.6), Inches(4.9), Inches(8.8), Inches(2.0), size=15, color=BLACK)
    add_footer(s, 5, total_slides)

    # ── Slide 6: Figure — Concentration and Beta ─────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Concentration and Beta by Condition",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_image(s, os.path.join(FIG_DIR, "fig1_hhi_by_condition.png"),
              Inches(0.3), Inches(1.4), width=Inches(4.7))
    add_image(s, os.path.join(FIG_DIR, "fig2_beta_by_condition.png"),
              Inches(5.0), Inches(1.4), width=Inches(4.7))
    add_body_box(s, [
        "Error bars are 95% CIs  •  all pairwise differences vs. control are p < 0.001",
    ], Inches(0.6), Inches(6.4), Inches(8.8), Inches(0.5), size=13, color=GREY)
    add_footer(s, 6, total_slides)

    # ── Slide 7: Figure — Sector Allocation ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Sector Allocation: The Story Is Sectoral",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_image(s, os.path.join(FIG_DIR, "fig4_sector_allocation_stacked.png"),
              Inches(1.5), Inches(1.3), width=Inches(7.0))
    add_body_box(s, [
        "Fundamental → Tech-heavy (NVDA, MSFT, GOOGL growth numbers drive it).",
        "Technical → more balanced; Energy and Healthcare boosted by MA signals.",
    ], Inches(0.6), Inches(6.4), Inches(8.8), Inches(0.5), size=13, color=GREY)
    add_footer(s, 7, total_slides)

    # ── Slide 8: Figure — Treatment Effects ──────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Treatment Effects vs. Control (95% CI)",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    add_image(s, os.path.join(FIG_DIR, "fig7_treatment_effects.png"),
              Inches(0.3), Inches(1.3), width=Inches(9.4))
    add_body_box(s, [
        "Every point lies well away from zero. Fundamental pushes metrics up; "
        "technical pulls beta and volatility down.",
    ], Inches(0.6), Inches(5.9), Inches(8.8), Inches(1.0), size=14, color=BLACK)
    add_footer(s, 8, total_slides)

    # ── Slide 9: Diagnostics ──────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Diagnostics — Are the Effects Real?",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))

    add_title_box(s, "Permutation placebo (5,000 shuffles)",
        Inches(0.6), Inches(1.4), Inches(8.8), Inches(0.5),
        size=17, bold=True, color=NAVY)
    add_body_box(s, [
        "Observed effects are outside the null distribution in every case.  "
        "Empirical p-value < 0.001 for 11 of 12 tests.",
    ], Inches(0.8), Inches(1.95), Inches(8.6), Inches(0.7), size=14)

    add_title_box(s, "Within-control placebo (1,000 random splits)",
        Inches(0.6), Inches(2.9), Inches(8.8), Inches(0.5),
        size=17, bold=True, color=NAVY)
    add_body_box(s, [
        "Rejection rates at α=0.05:  {HHI 0.058,  sector HHI 0.047,  "
        "β 0.057,  vol 0.049}.   Calibrated to nominal 5%.",
    ], Inches(0.8), Inches(3.45), Inches(8.6), Inches(0.9), size=14)

    add_title_box(s, "Trimmed-sample sensitivity",
        Inches(0.6), Inches(4.65), Inches(8.8), Inches(0.5),
        size=17, bold=True, color=NAVY)
    add_body_box(s, [
        "Drop the top and bottom HHI observation per condition (n: 119 → 111).  "
        "All point estimates change by <5%.  Every effect remains p < 0.001.",
    ], Inches(0.8), Inches(5.2), Inches(8.6), Inches(0.9), size=14)

    add_title_box(s,
        "The effects are not driven by noise, pipeline miscalibration, or outliers.",
        Inches(0.6), Inches(6.4), Inches(8.8), Inches(0.5),
        size=14, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_footer(s, 9, total_slides)

    # ── Slide 10: Discussion / Conclusion ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    add_title_box(s, "Discussion and Conclusion",
        Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))

    add_title_box(s, "Takeaways",
        Inches(0.6), Inches(1.4), Inches(4.5), Inches(0.5),
        size=18, bold=True, color=NAVY)
    add_body_box(s, [
        "Information framing is a causal driver of LLM portfolio behavior.",
        "The combined prompt is NOT a superposition — it de-concentrates.",
        "The tools of randomized experiments & placebos carry over cleanly to LLMs.",
    ], Inches(0.6), Inches(1.95), Inches(8.8), Inches(2.4), size=15)

    add_title_box(s, "Limitations",
        Inches(0.6), Inches(4.35), Inches(4.5), Inches(0.5),
        size=18, bold=True, color=NAVY)
    add_body_box(s, [
        "One model (Claude Sonnet 4.6), one temperature, one snapshot of data.",
        "Prompts differ in length as well as content — no placebo arm for length.",
        "Breadth is a binding constraint; cannot vary across conditions by design.",
        "No out-of-sample evaluation: studies what the LLM does, not whether it's good.",
    ], Inches(0.6), Inches(4.9), Inches(8.8), Inches(2.0), size=14)

    add_title_box(s,
        "Code & data: github.com/eddieschwasnick/CI-Final-Project",
        Inches(0.6), Inches(6.9), Inches(8.8), Inches(0.4),
        size=14, bold=True, color=GREY, align=PP_ALIGN.CENTER)
    add_footer(s, 10, total_slides)

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH} ({len(prs.slides)} slides).")


if __name__ == "__main__":
    main()
